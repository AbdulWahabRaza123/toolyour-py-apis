"""
Office conversion service.
Handles XLS, XLSX, PPT, PPTX, ODT, ODS, ODP conversions.
"""

import io
import csv
from typing import Optional
import structlog
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import xlrd
import xlwt

from .types import OfficeServiceResponse, OfficeConversionOptions

logger = structlog.get_logger(__name__)


class OfficeConverterService:
    """Service for converting office formats."""

    def __init__(self):
        self.supported_conversions = {
            'xls': ['xlsx', 'csv', 'txt', 'json'],
            'xlsx': ['xls', 'csv', 'txt', 'json', 'pdf'],
            'ppt': ['pptx', 'pdf', 'txt', 'html'],
            'pptx': ['ppt', 'pdf', 'txt', 'html', 'images'],
            'odt': ['docx', 'txt', 'pdf', 'html'],
            'ods': ['xlsx', 'csv', 'txt', 'json'],
            'odp': ['pptx', 'pdf', 'txt', 'html'],
        }

    def can_convert(self, source_format: str, target_format: str) -> bool:
        """Check if conversion is supported."""
        source_format = source_format.lower().replace('.', '')
        target_format = target_format.lower().replace('.', '')
        
        return target_format in self.supported_conversions.get(source_format, [])

    def get_supported_formats(self, source_format: str) -> list:
        """Get supported target formats for a source format."""
        source_format = source_format.lower().replace('.', '')
        return self.supported_conversions.get(source_format, [])

    # XLS conversions
    async def convert_xls_to_xlsx(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLS to XLSX."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLS file
            xls_workbook = xlrd.open_workbook(file_contents=file_buffer)
            
            # Create new XLSX workbook
            xlsx_workbook = Workbook()
            xlsx_workbook.remove(xlsx_workbook.active)  # Remove default sheet
            
            for sheet_name in xls_workbook.sheet_names():
                xls_sheet = xls_workbook.sheet_by_name(sheet_name)
                xlsx_sheet = xlsx_workbook.create_sheet(title=sheet_name)
                
                # Copy data
                for row in range(xls_sheet.nrows):
                    for col in range(xls_sheet.ncols):
                        cell_value = xls_sheet.cell_value(row, col)
                        xlsx_sheet.cell(row=row+1, column=col+1, value=cell_value)

            # Save to bytes
            xlsx_buffer = io.BytesIO()
            xlsx_workbook.save(xlsx_buffer)
            xlsx_content = xlsx_buffer.getvalue()
            xlsx_buffer.close()

            logger.info("XLS to XLSX conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLS converted to XLSX successfully",
                data=xlsx_content,
                format="xlsx"
            )

        except Exception as e:
            logger.error("XLS to XLSX conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLS to XLSX",
                error=str(e)
            )

    async def convert_xls_to_csv(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLS to CSV."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLS file
            xls_workbook = xlrd.open_workbook(file_contents=file_buffer)
            sheet_name = options.sheet_name or xls_workbook.sheet_names()[0]
            xls_sheet = xls_workbook.sheet_by_name(sheet_name)
            
            # Convert to CSV
            csv_buffer = io.StringIO()
            writer = csv.writer(csv_buffer)
            
            for row in range(xls_sheet.nrows):
                row_data = []
                for col in range(xls_sheet.ncols):
                    cell_value = xls_sheet.cell_value(row, col)
                    row_data.append(str(cell_value))
                writer.writerow(row_data)
            
            csv_content = csv_buffer.getvalue()
            csv_buffer.close()

            logger.info("XLS to CSV conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLS converted to CSV successfully",
                data=csv_content.encode(options.encoding),
                format="csv"
            )

        except Exception as e:
            logger.error("XLS to CSV conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLS to CSV",
                error=str(e)
            )

    async def convert_xls_to_txt(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLS to TXT."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLS file
            xls_workbook = xlrd.open_workbook(file_contents=file_buffer)
            sheet_name = options.sheet_name or xls_workbook.sheet_names()[0]
            xls_sheet = xls_workbook.sheet_by_name(sheet_name)
            
            # Convert to text
            text_lines = []
            for row in range(xls_sheet.nrows):
                row_data = []
                for col in range(xls_sheet.ncols):
                    cell_value = xls_sheet.cell_value(row, col)
                    row_data.append(str(cell_value))
                text_lines.append("\t".join(row_data))
            
            text_content = "\n".join(text_lines)

            logger.info("XLS to TXT conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLS converted to TXT successfully",
                data=text_content.encode(options.encoding),
                format="txt"
            )

        except Exception as e:
            logger.error("XLS to TXT conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLS to TXT",
                error=str(e)
            )

    async def convert_xls_to_json(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLS to JSON."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLS file
            xls_workbook = xlrd.open_workbook(file_contents=file_buffer)
            sheet_name = options.sheet_name or xls_workbook.sheet_names()[0]
            xls_sheet = xls_workbook.sheet_by_name(sheet_name)
            
            # Convert to JSON
            json_data = {
                "sheet_name": sheet_name,
                "data": []
            }
            
            for row in range(xls_sheet.nrows):
                row_data = {}
                for col in range(xls_sheet.ncols):
                    cell_value = xls_sheet.cell_value(row, col)
                    row_data[f"col_{col}"] = cell_value
                json_data["data"].append(row_data)
            
            import json
            json_str = json.dumps(json_data, indent=2)

            logger.info("XLS to JSON conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLS converted to JSON successfully",
                data=json_str.encode(options.encoding),
                format="json"
            )

        except Exception as e:
            logger.error("XLS to JSON conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLS to JSON",
                error=str(e)
            )

    # XLSX conversions
    async def convert_xlsx_to_xls(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLSX to XLS."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLSX file
            xlsx_workbook = load_workbook(io.BytesIO(file_buffer))
            
            # Create new XLS workbook
            xls_workbook = xlwt.Workbook()
            
            for sheet_name in xlsx_workbook.sheetnames:
                xlsx_sheet = xlsx_workbook[sheet_name]
                xls_sheet = xls_workbook.add_sheet(sheet_name)
                
                # Copy data
                for row in xlsx_sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            xls_sheet.write(cell.row-1, cell.column-1, cell.value)

            # Save to bytes
            xls_buffer = io.BytesIO()
            xls_workbook.save(xls_buffer)
            xls_content = xls_buffer.getvalue()
            xls_buffer.close()

            logger.info("XLSX to XLS conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLSX converted to XLS successfully",
                data=xls_content,
                format="xls"
            )

        except Exception as e:
            logger.error("XLSX to XLS conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLSX to XLS",
                error=str(e)
            )

    async def convert_xlsx_to_csv(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLSX to CSV."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLSX file
            xlsx_workbook = load_workbook(io.BytesIO(file_buffer))
            sheet_name = options.sheet_name or xlsx_workbook.sheetnames[0]
            xlsx_sheet = xlsx_workbook[sheet_name]
            
            # Convert to CSV
            csv_buffer = io.StringIO()
            writer = csv.writer(csv_buffer)
            
            for row in xlsx_sheet.iter_rows(values_only=True):
                writer.writerow(row)
            
            csv_content = csv_buffer.getvalue()
            csv_buffer.close()

            logger.info("XLSX to CSV conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLSX converted to CSV successfully",
                data=csv_content.encode(options.encoding),
                format="csv"
            )

        except Exception as e:
            logger.error("XLSX to CSV conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLSX to CSV",
                error=str(e)
            )

    async def convert_xlsx_to_txt(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLSX to TXT."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLSX file
            xlsx_workbook = load_workbook(io.BytesIO(file_buffer))
            sheet_name = options.sheet_name or xlsx_workbook.sheetnames[0]
            xlsx_sheet = xlsx_workbook[sheet_name]
            
            # Convert to text
            text_lines = []
            for row in xlsx_sheet.iter_rows(values_only=True):
                text_lines.append("\t".join(str(cell) if cell is not None else "" for cell in row))
            
            text_content = "\n".join(text_lines)

            logger.info("XLSX to TXT conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLSX converted to TXT successfully",
                data=text_content.encode(options.encoding),
                format="txt"
            )

        except Exception as e:
            logger.error("XLSX to TXT conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLSX to TXT",
                error=str(e)
            )

    async def convert_xlsx_to_json(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert XLSX to JSON."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read XLSX file
            xlsx_workbook = load_workbook(io.BytesIO(file_buffer))
            sheet_name = options.sheet_name or xlsx_workbook.sheetnames[0]
            xlsx_sheet = xlsx_workbook[sheet_name]
            
            # Convert to JSON
            json_data = {
                "sheet_name": sheet_name,
                "data": []
            }
            
            for row in xlsx_sheet.iter_rows(values_only=True):
                row_data = {}
                for i, cell in enumerate(row):
                    row_data[f"col_{i}"] = cell
                json_data["data"].append(row_data)
            
            import json
            json_str = json.dumps(json_data, indent=2)

            logger.info("XLSX to JSON conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="XLSX converted to JSON successfully",
                data=json_str.encode(options.encoding),
                format="json"
            )

        except Exception as e:
            logger.error("XLSX to JSON conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting XLSX to JSON",
                error=str(e)
            )

    # PPT conversions
    async def convert_ppt_to_pptx(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert PPT to PPTX."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # For now, return a placeholder response
            # PPT to PPTX conversion requires more complex handling
            logger.warning("PPT to PPTX conversion not fully implemented")
            return OfficeServiceResponse(
                status=501,
                message="PPT to PPTX conversion requires additional libraries",
                error="PPT to PPTX conversion not implemented"
            )

        except Exception as e:
            logger.error("PPT to PPTX conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting PPT to PPTX",
                error=str(e)
            )

    async def convert_ppt_to_txt(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert PPT to TXT."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # For now, return a placeholder response
            logger.warning("PPT to TXT conversion not fully implemented")
            return OfficeServiceResponse(
                status=501,
                message="PPT to TXT conversion requires additional libraries",
                error="PPT to TXT conversion not implemented"
            )

        except Exception as e:
            logger.error("PPT to TXT conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting PPT to TXT",
                error=str(e)
            )

    # PPTX conversions
    async def convert_pptx_to_txt(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert PPTX to TXT."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read PPTX file
            presentation = Presentation(io.BytesIO(file_buffer))
            
            # Extract text from slides
            text_content = []
            for i, slide in enumerate(presentation.slides):
                if options.slide_number is None or i == options.slide_number - 1:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        text_content.append(f"Slide {i+1}:")
                        text_content.append("\n".join(slide_text))
                        text_content.append("")

            result_text = "\n".join(text_content)

            logger.info("PPTX to TXT conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="PPTX converted to TXT successfully",
                data=result_text.encode(options.encoding),
                format="txt"
            )

        except Exception as e:
            logger.error("PPTX to TXT conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting PPTX to TXT",
                error=str(e)
            )

    async def convert_pptx_to_html(
        self,
        file_buffer: bytes,
        options: Optional[OfficeConversionOptions] = None
    ) -> OfficeServiceResponse:
        """Convert PPTX to HTML."""
        try:
            if options is None:
                options = OfficeConversionOptions()

            # Read PPTX file
            presentation = Presentation(io.BytesIO(file_buffer))
            
            # Convert to HTML
            html_content = f'<!DOCTYPE html>\n<html>\n<head>\n<meta charset="{options.encoding}">\n<title>Presentation</title>\n</head>\n<body>\n'
            
            for i, slide in enumerate(presentation.slides):
                if options.slide_number is None or i == options.slide_number - 1:
                    html_content += f'<div class="slide">\n<h2>Slide {i+1}</h2>\n'
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            html_content += f'<p>{shape.text}</p>\n'
                    
                    html_content += '</div>\n'
            
            html_content += '</body>\n</html>'

            logger.info("PPTX to HTML conversion completed")
            return OfficeServiceResponse(
                status=200,
                message="PPTX converted to HTML successfully",
                data=html_content.encode(options.encoding),
                format="html"
            )

        except Exception as e:
            logger.error("PPTX to HTML conversion failed", error=str(e))
            return OfficeServiceResponse(
                status=500,
                message="Error converting PPTX to HTML",
                error=str(e)
            )

    async def get_supported_conversions(self):
        """Get list of supported office conversions."""
        return {
            "supported_conversions": self.supported_conversions,
            "message": "List of supported office format conversions"
        }


# Global instance
office_converter_service = OfficeConverterService()
